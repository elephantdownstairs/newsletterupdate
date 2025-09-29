import streamlit as st
import pandas as pd
import io
import zipfile
import re
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.text import PP_ALIGN  # Make sure this import is present at the top
import openpyxl

def sanitize_filename(filename):
    """Sanitize filename for safe saving"""
    return re.sub(r'[<>:"/\\|?*]', '_', filename)

def format_date(date_str):
    """Format date from YYYY-MM-DD to 'DD MMM YYYY'"""
    if not date_str or pd.isna(date_str):
        return ""
    
    try:
        if isinstance(date_str, str):
            date_obj = datetime.strptime(date_str, "%Y-%m-%d")
        else:
            date_obj = date_str
        return date_obj.strftime("%d %b %Y")
    except:
        return str(date_str)

# ========================================
# 🔧 CHANGE 1: UPDATED PERCENTAGE FORMATTING
# ========================================
def format_percentage(value):
    """Format numeric values as percentages with robust handling"""
    if pd.isna(value) or value == "":
        return ""
    
    try:
        # Handle string input
        if isinstance(value, str):
            value = value.strip()
            if value == '':
                return ""
            
            # Check if already has % symbol
            if '%' in value:
                # Remove % and convert to number
                num_value = float(value.replace('%', '').strip())
                # Already in percentage form (e.g., "2%" means 2.00%)
                return f"{num_value:.2f}%"
            else:
                # String number without %
                num_value = float(value)
                # If <= 1, treat as decimal (0.02 → 2.00%)
                # If > 1, treat as already percentage (2 → 2.00%)
                if num_value <= 1:
                    return f"{num_value * 100:.2f}%"
                else:
                    return f"{num_value:.2f}%"
        else:
            # Numeric input
            num_value = float(value)
            # If <= 1, treat as decimal (0.02 → 2.00%)
            # If > 1, treat as already percentage (2 → 2.00%)
            if num_value <= 1:
                return f"{num_value * 100:.2f}%"
            else:
                return f"{num_value:.2f}%"
    except:
        return str(value)
# ========================================

def is_percentage_field(placeholder):
    """Check if a placeholder should be formatted as percentage"""
    percentage_fields = [
        'A_1M', 'B_1M', 'C_1M',
        'A_YTD', 'B_YTD', 'C_YTD', 
        'A_1Y', 'B_1Y', 'C_1Y',
        'A_SI', 'B_SI', 'C_SI',
        'A_LastY', 'B_LastY', 'C_LastY'
    ]
    return placeholder in percentage_fields

def get_font_color(placeholder):
    """Determine font color based on placeholder"""
    white_placeholders = ['fund_name', 'slide_date']
    if placeholder in white_placeholders:
        return RGBColor(255, 255, 255)  # White
    else:
        return RGBColor(0, 0, 0)  # Black

def get_font_size(placeholder):
    """Determine font size based on placeholder"""
    font_sizes = {
        'fund_name': 22.5,
        'slide_date': 15,
        'investment_objective': 9,
        'performance_review': 9,
        'equity_strategy': 9,
        'fixed_income_strategy': 9,
        'table_note': 7,
        'Disclaimer': 7,
        'date': 7,
        'class_A': 10,
        'class_B': 10,
        'class_C': 10,
        'A_1M': 10, 'B_1M': 10, 'C_1M': 10,
        'A_YTD': 10, 'B_YTD': 10, 'C_YTD': 10,
        'A_1Y': 10, 'B_1Y': 10, 'C_1Y': 10,
        'A_SI': 10, 'B_SI': 10, 'C_SI': 10,
        'A_LastY': 10, 'B_LastY': 10, 'C_LastY': 10
    }
    return font_sizes.get(placeholder, 10)  # Default to 10 if not specified

def should_be_bold(placeholder):
    """Determine if placeholder should be bold"""
    bold_placeholders = ['fund_name', 'class_A', 'class_B', 'class_C']
    return placeholder in bold_placeholders

# List of placeholders that should be center-aligned
CENTER_ALIGN_PLACEHOLDERS = [
    'class_A', 'class_B', 'class_C',
    'A_1M', 'B_1M', 'C_1M',
    'A_YTD', 'B_YTD', 'C_YTD',
    'A_1Y', 'B_1Y', 'C_1Y',
    'A_SI', 'B_SI', 'C_SI',
    'A_LastY', 'B_LastY', 'C_LastY'
]

def replace_text_in_shape(shape, replacements):
    """Replace placeholders in a shape's text and apply formatting"""
    if not shape.has_text_frame:
        return

    text_frame = shape.text_frame

    for paragraph in text_frame.paragraphs:
        full_text = paragraph.text
        placeholders = re.findall(r'\{\{\s*([^}]+)\s*\}\}', full_text)

        if placeholders:
            new_text = full_text
            for placeholder in placeholders:
                placeholder_key = placeholder.strip()

                if placeholder_key in replacements:
                    if is_percentage_field(placeholder_key):
                        replacement_value = format_percentage(replacements[placeholder_key])
                    elif placeholder_key == 'date':
                        replacement_value = format_date(replacements[placeholder_key])
                    else:
                        replacement_value = replacements[placeholder_key]
                else:
                    replacement_value = ""

                new_text = new_text.replace(f"{{{{{placeholder}}}}}", str(replacement_value))

            paragraph.clear()
            run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
            run.text = new_text

            for placeholder in placeholders:
                placeholder_key = placeholder.strip()
                if placeholder_key in replacements:
                    run.font.name = "Poppins"
                    run.font.color.rgb = get_font_color(placeholder_key)
                    run.font.size = Pt(get_font_size(placeholder_key))
                    run.font.bold = should_be_bold(placeholder_key)

                    # ✅ Align center if in defined list
                    if placeholder_key in CENTER_ALIGN_PLACEHOLDERS:
                        paragraph.alignment = PP_ALIGN.CENTER
                    break


# ========================================
# 🔧 CHANGE 2: UPDATED TABLE CELL FORMATTING
# Added center alignment and vertical middle anchor
# ========================================
def replace_text_in_table(table, replacements):
    """Replace placeholders in table cells with formatting and alignment"""
    for row in table.rows:
        for cell in row.cells:
            if cell.text:
                cell_text = cell.text
                placeholders = re.findall(r'\{\{\s*([^}]+)\s*\}\}', cell_text)
                
                if placeholders:
                    new_text = cell_text
                    font_color = None
                    font_size = None
                    is_bold = False
                    should_center = False  # 🆕 NEW: Track if cell should be centered
                    
                    for placeholder in placeholders:
                        placeholder_key = placeholder.strip()
                        
                        # Get replacement value with proper formatting
                        if placeholder_key in replacements:
                            if is_percentage_field(placeholder_key):
                                replacement_value = format_percentage(replacements[placeholder_key])
                            elif placeholder_key == 'date':
                                replacement_value = format_date(replacements[placeholder_key])
                            else:
                                replacement_value = replacements[placeholder_key]
                        else:
                            replacement_value = ""
                        
                        new_text = new_text.replace(f"{{{{{placeholder}}}}}", str(replacement_value))
                        
                        # Set font properties based on first placeholder found
                        if font_color is None:
                            font_color = get_font_color(placeholder_key)
                            font_size = get_font_size(placeholder_key)
                            is_bold = should_be_bold(placeholder_key)
                            should_center = placeholder_key in CENTER_ALIGN_PLACEHOLDERS  # 🆕 NEW
                    
                    # Update cell text
                    cell.text = new_text
                    
                    # Apply font formatting to all paragraphs in the cell
                    for paragraph in cell.text_frame.paragraphs:
                        # 🆕 NEW: Center align if needed
                        if should_center:
                            paragraph.alignment = PP_ALIGN.CENTER
                        
                        for run in paragraph.runs:
                            run.font.name = "Poppins"
                            if font_color:
                                run.font.color.rgb = font_color
                            if font_size:
                                run.font.size = Pt(font_size)
                            run.font.bold = is_bold
                    
                    # 🆕 NEW: Set vertical anchor to middle for center-aligned cells
                    if should_center:
                        cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
# ========================================

def should_show_section(row, section_name):
    """Check if a section should be shown based on show_ column"""
    show_column = f"show_{section_name}"
    if show_column in row.index:
        show_value = row[show_column]
        if pd.isna(show_value) or str(show_value).strip().lower() in ['no', 'false', '0', '']:
            return False
    return True

def process_slide(slide, row):
    """Process a single slide with fund data"""
    # Prepare replacements dictionary
    replacements = {}
    
    for column in row.index:
        value = row[column]
        if pd.isna(value):
            replacements[column] = ""
        elif column == 'slide_date':
            replacements[column] = format_date(value)
        elif column == 'date':
            replacements[column] = format_date(value)
        else:
            replacements[column] = str(value)
    
    # Process all shapes in the slide
    for shape in slide.shapes:
        if shape.has_text_frame:
            replace_text_in_shape(shape, replacements)
        elif shape.has_table:
            replace_text_in_table(shape.table, replacements)
    
    # Handle conditional sections
    sections_to_check = ['equity_strategy', 'fixed_income_strategy']
    
    for section in sections_to_check:
        if not should_show_section(row, section):
            # Find and remove or clear shapes with this section's placeholder
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    if f"{{{{{section}}}}}" in text:
                        shape.text_frame.clear()

def generate_presentation(template_pptx, fund_data):
    """Generate a presentation for a single fund"""
    presentation = Presentation(template_pptx)
    
    # Process each slide
    for slide in presentation.slides:
        process_slide(slide, fund_data)
    
    return presentation

def main():
    st.set_page_config(
        page_title="Fund Update Generator",
        page_icon="📊",
        layout="wide"
    )
    
    st.title("📊 PowerPoint Fund Update Generator")
    st.markdown("Upload your Excel data and PowerPoint template to generate automated fund updates.")
    
    # File upload section
    col1, col2 = st.columns(2)
    
    with col1:
        st.subheader("📁 Upload Excel Template")
        excel_file = st.file_uploader(
            "Choose Excel file",
            type=['xlsx', 'xls'],
            help="Upload Fund_Update_Template_v1.xlsx"
        )
    
    with col2:
        st.subheader("📁 Upload PowerPoint Template")
        pptx_file = st.file_uploader(
            "Choose PowerPoint file",
            type=['pptx'],
            help="Upload MAM-Fund Update Template PPT.pptx"
        )
    
    if excel_file and pptx_file:
        try:
            # Read Excel file
            df = pd.read_excel(excel_file)
            st.success(f"✅ Excel file loaded successfully! Found {len(df)} funds.")
            
            # Display fund preview
            st.subheader("📋 Fund Preview")
            if 'fund_name' in df.columns:
                fund_names = df['fund_name'].tolist()
                st.write(f"**Funds found:** {', '.join(map(str, fund_names))}")
                
                # Fund selection
                st.subheader("🎯 Select Funds to Generate")
                selected_funds = st.multiselect(
                    "Choose funds to generate presentations for:",
                    options=fund_names,
                    default=fund_names[:min(3, len(fund_names))]  # Default to first 3 funds
                )
                
                if selected_funds:
                    # Generation options
                    col1, col2 = st.columns(2)
                    with col1:
                        generate_individual = st.button("🔄 Generate Individual Files", type="primary")
                    with col2:
                        generate_zip = st.button("📦 Generate ZIP Package")
                    
                    if generate_individual or generate_zip:
                        progress_bar = st.progress(0)
                        status_text = st.empty()
                        
                        generated_files = []
                        
                        for i, fund_name in enumerate(selected_funds):
                            status_text.text(f"Processing {fund_name}...")
                            
                            # Get fund data
                            fund_row = df[df['fund_name'] == fund_name].iloc[0]
                            
                            # Generate presentation
                            presentation = generate_presentation(pptx_file, fund_row)
                            
                            # Save to memory
                            output = io.BytesIO()
                            presentation.save(output)
                            output.seek(0)
                            
                            # Create filename
                            safe_fund_name = sanitize_filename(str(fund_name))
                            filename = f"{safe_fund_name}_Fund_Update.pptx"
                            
                            generated_files.append({
                                'name': filename,
                                'data': output.getvalue(),
                                'fund_name': fund_name
                            })
                            
                            progress_bar.progress((i + 1) / len(selected_funds))
                        
                        status_text.text("✅ Generation complete!")
                        
                        if generate_individual:
                            st.subheader("📥 Download Individual Files")
                            for file_info in generated_files:
                                st.download_button(
                                    label=f"📄 Download {file_info['fund_name']}",
                                    data=file_info['data'],
                                    file_name=file_info['name'],
                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                                )
                        
                        if generate_zip:
                            # Create ZIP file
                            zip_buffer = io.BytesIO()
                            with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                                for file_info in generated_files:
                                    zip_file.writestr(file_info['name'], file_info['data'])
                            
                            zip_buffer.seek(0)
                            
                            st.subheader("📦 Download ZIP Package")
                            st.download_button(
                                label="📦 Download All Files (ZIP)",
                                data=zip_buffer.getvalue(),
                                file_name=f"Fund_Updates_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip",
                                mime="application/zip"
                            )
            else:
                st.error("❌ Excel file must contain a 'fund_name' column")
                
        except Exception as e:
            st.error(f"❌ Error processing files: {str(e)}")
            st.exception(e)
    
    elif excel_file or pptx_file:
        st.info("📋 Please upload both Excel and PowerPoint template files to continue.")
    
    else:
        st.info("👆 Please upload your Excel data file and PowerPoint template to get started.")
    
    # Instructions
    with st.expander("📖 Instructions"):
        st.markdown("""
        ### How to use this tool:
        
        1. **Upload Files**: 
           - Upload your Excel template (Fund_Update_Template_v1.xlsx)
           - Upload your PowerPoint template (MAM-Fund Update Template PPT.pptx)
        
        2. **Select Funds**: Choose which funds you want to generate presentations for
        
        3. **Generate**: Click either button to generate individual files or a ZIP package
        
        ### Excel Template Requirements:
        - Must contain a 'fund_name' column
        - Should include columns matching PowerPoint placeholders (e.g., slide_date, investment_objective)
        - Use 'show_equity_strategy' and 'show_fixed_income_strategy' columns to control section visibility
        
        ### PowerPoint Template Requirements:
        - Use placeholders in format `{{ placeholder_name }}`
        - Tables should have tagged cells (e.g., `{{A_1M}}`, `{{B_YTD}}`)
        - fund_name and slide_date will be colored white, others black
        
        ### Features:
        - ✅ Automatic placeholder replacement
        - ✅ Font color formatting (white for fund_name/slide_date, black for others)
        - ✅ Date formatting (YYYY-MM-DD → "DD MMM YYYY")
        - ✅ Robust percentage formatting (0.02 → 2.00%, "2%" → 2.00%)
        - ✅ Center-aligned table cells with middle vertical anchor
        - ✅ Poppins font family applied to all text
        - ✅ Conditional section handling
        - ✅ Table cell replacement
        - ✅ Safe filename generation
        - ✅ Individual or batch download
        """)

if __name__ == "__main__":
    main()
